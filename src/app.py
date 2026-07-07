import sys
import os
import streamlit as st
from dotenv import load_dotenv
import pandas as pd
from pypdf import PdfReader
from pptx import Presentation

# --- Path Injection Fix ---
root_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if root_path not in sys.path:
    sys.path.append(root_path)

from src.main import run_agent

# --- Setup & Config ---
load_dotenv()
st.set_page_config(page_title="Mediator & Logic Engine", layout="wide", page_icon="🛡️")

# --- Helper: Text Extraction ---
def get_file_text(file):
    file.seek(0)
    text = f"\n--- USER DOCUMENT: {file.name} ---\n"
    try:
        if file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif "excel" in file.type or "spreadsheet" in file.type:
            df = pd.read_excel(file)
            text += df.to_string()
        elif "powerpoint" in file.type or "presentation" in file.type:
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            text += file.getvalue().decode("utf-8", errors="ignore")
    except Exception as e:
        text += f"[Error reading file: {e}]"
    return text

# --- Session State ---
if "messages" not in st.session_state:
    st.session_state.messages = []

def main():
    st.markdown("<h1>🛡️ Mediator & Logic Engine</h1>", unsafe_allow_html=True)
    st.caption("Autonomous Policy Arbitration & Universal Reasoning System")

    with st.sidebar:
        st.markdown("## ⚙️ CONTROL ARRAY")
        st.write("📦 **Pre-loaded Knowledge:**")
        st.caption("✅ `policy.txt` (vector-indexed)")
        st.caption("✅ `synthetic_learners.json`")
        st.divider()
        st.caption("_Note: pre-loaded policy & performance data are fictional samples, not real company data._")

        uploaded_files = st.file_uploader(
            "Upload Custom Docs",
            accept_multiple_files=True,
            type=['pdf', 'xlsx', 'xls', 'pptx', 'ppt', 'txt'],
            key="my_uploader"
        )
        st.divider()
        st.success("Foundry/Fabric/Work IQ Integrated")
        if st.button("Clear Conversation"):
            st.session_state.messages = []
            st.rerun()

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    if user_query := st.chat_input("Ask about company policy, team performance, or general research..."):
        st.session_state.messages.append({"role": "user", "content": user_query})
        with st.chat_message("user"):
            st.markdown(user_query)

        with st.chat_message("assistant"):
            with st.spinner("Agent is reasoning..."):
                try:
                    uploaded_text = ""
                    if uploaded_files:
                        for file in uploaded_files:
                            uploaded_text += get_file_text(file) + "\n"
                        st.info(f"✅ Including {len(uploaded_files)} file(s) in analysis.")

                    result = run_agent(user_query, uploaded_text=uploaded_text)

                    with st.expander("🔍 Reasoning Trace & Observability"):
                        st.write(f"**Targeted Path:** {result.get('path', 'Unknown').upper()}")
                        st.write(f"**Orchestration Status:** Success")

                    st.markdown(result['content'])
                    st.session_state.messages.append({"role": "assistant", "content": result['content']})
                except Exception as e:
                    error_msg = f"⚠️ Execution Error: {e}"
                    st.error(error_msg)
                    st.session_state.messages.append({"role": "assistant", "content": error_msg})
            st.rerun()

if __name__ == "__main__":
    main()